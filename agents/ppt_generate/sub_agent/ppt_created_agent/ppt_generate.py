from pptx import Presentation
import re

class ppt_generate_class:
    def __init__(self, no_of_slides, topic_name, title, content):
        self.no_of_slides = no_of_slides
        self.topic_name = topic_name
        self.title = title
        self.content = content

    def generate_ppt(self):
        prs = Presentation()
        print(f"Generating PPT with {self.no_of_slides} slides on topic: {self.topic_name}...")
        safe_topic = re.sub(r'[^\w\-]', '_', self.topic_name)  # Replace invalid chars
        for i in range(self.no_of_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = self.title[i]
            content.text = self.content[i]
        prs.save(f"{safe_topic}.pptx")